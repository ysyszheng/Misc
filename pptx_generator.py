import collections 
import collections.abc
import re
import time
from pptx import Presentation
from pptx.dml.color import RGBColor
# from pptx.util import Cm, Pt

def move_slide(presentation, old_index, new_index):
        xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])

file_md_path = r"./question.md"
file_pptx_path = r"./question.pptx"
total = 224
cnt = 0
time_begin = time.time()

with open(file_md_path, "r", encoding='UTF-8') as file:
    # patternQ = re.compile(r'^\s*([1-9][0-9]*\. )(.*)(==)(.*)(==)(.*)$') # markdown
    # patternQ = re.compile(r'^\s*([1-9][0-9]*[\.、]?\s*)(.*)(（|==|【)(.*)(）|==|】)(.*)$')
    patternQ = re.compile(r'^\s*([1-9][0-9]*[\.、]?\s*)(.*)(【)(.*)(】)(.*)$')
    patternA = re.compile(r'^\s*\D.*$')
    bk_line = ''

    while True:
        if bk_line:
            line = bk_line
        else:
            line = file.readline()
        
        if line and re.match(patternQ, line):
            question = re.sub(patternQ, r'\2【 】\6', line)
            solution = re.sub(patternQ, r'\2【\4】\6', line)
            qleft = re.sub(patternQ, r'\2', line)
            qright = re.sub(patternQ, r'\6', line)
            amid = re.sub(patternQ, r'【\4】', line)
            ans = re.sub(patternQ, r'\4', line)

            qleft = qleft[:-1]
            amid = amid[:-1]

            prs = Presentation(file_pptx_path)
            slide_layout = prs.slide_layouts[3]
            slide = prs.slides.add_slide(slide_layout)

            textbox = slide.placeholders[10]
            tf = textbox.text_frame
            tf.clear()

            para = tf.paragraphs[0]
            qleft_run = para.add_run()
            qleft_run.text = qleft
            amid_run = para.add_run()
            amid_run.text = amid
            amid_run.font.color.rgb = RGBColor(255, 217, 102)
            amid_run.font.bold = True
            qright_run = para.add_run()
            qright_run.text = qright
            
            while True:
                newline = file.readline()
                if newline and re.match(patternA, newline):
                    question += newline
                    solution += newline

                    if newline[0] in ans:
                        new_run = para.add_run()
                        new_run.text = newline
                        font = new_run.font
                        font.color.rgb = RGBColor(255, 217, 102)
                        font.bold = True
                    else:
                        new_run = para.add_run()
                        new_run.text = newline
                        font = new_run.font
                        font.color.rgb = RGBColor(255, 255, 255)
                else:
                    bk_line = newline 
                    break
            
            prs.save(file_pptx_path)

            cnt += 1
            time_end = time.time()
            print(solution)
            print('-----------', cnt/total*100, '%, time cost', time_end-time_begin, 's -----------')

            prs = Presentation(file_pptx_path)
            title_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(title_slide_layout)
            content = slide.placeholders[10]
            content.text = question

            move_slide(prs, prs.slides.index(slide), prs.slides.index(slide)-1)

            prs.save(file_pptx_path)

        if line == '':
            break

